import re
from pathlib import Path
from typing import List, Union

from langchain_core.documents import Document
from langchain_community.document_loaders import Docx2txtLoader, PyMuPDFLoader, TextLoader


PathLike = Union[str, Path]

SUPPORTED_EXTENSIONS = {".pdf", ".txt", ".docx", ".pptx"}


def clean_text(text: str) -> str:
    text = text.replace("\u00a0", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _load_pptx_documents(file_path: Path) -> List[Document]:
    from pptx import Presentation

    presentation = Presentation(str(file_path))
    documents: List[Document] = []

    for slide_index, slide in enumerate(presentation.slides):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text.strip())

        page_content = "\n".join(parts).strip()
        if not page_content:
            continue

        documents.append(
            Document(
                page_content=page_content,
                metadata={"page": slide_index},
            )
        )

    return documents


def _load_raw_documents(file_path: Path) -> List[Document]:
    ext = file_path.suffix.lower()

    if ext == ".pdf":
        return PyMuPDFLoader(str(file_path)).load()
    if ext == ".txt":
        return TextLoader(str(file_path), encoding="utf-8").load()
    if ext == ".docx":
        return Docx2txtLoader(str(file_path)).load()
    if ext == ".pptx":
        return _load_pptx_documents(file_path)

    return []


def _attach_metadata(documents: List[Document], file_path: Path) -> List[Document]:
    subject = file_path.parent.name
    file_name = file_path.name

    for doc in documents:
        doc.page_content = clean_text(doc.page_content)
        doc.metadata["subject"] = subject
        doc.metadata["file_name"] = file_name
        doc.metadata["source"] = file_name
        doc.metadata["file_path"] = str(file_path)

        if "page" not in doc.metadata:
            doc.metadata["page"] = None

    return documents


def _load_file(file_path: Path) -> List[Document]:
    ext = file_path.suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        return []

    loaded_docs = _load_raw_documents(file_path)
    if not loaded_docs:
        return []

    return _attach_metadata(loaded_docs, file_path)


def load_documents(data_dir: str) -> List[Document]:
    """
    data_dir 폴더 안의 PDF/TXT/DOCX/PPTX 파일을 LangChain Document 객체로 로드합니다.
    """
    documents: List[Document] = []
    data_path = Path(data_dir)

    if not data_path.exists():
        raise FileNotFoundError(f"데이터 폴더를 찾을 수 없습니다: {data_dir}")

    for file_path in data_path.rglob("*"):
        if not file_path.is_file():
            continue

        documents.extend(_load_file(file_path))

    return documents


def load_documents_from_paths(file_paths: List[PathLike]) -> List[Document]:
    """지정한 PDF/TXT/DOCX/PPTX 파일만 Document로 로드합니다."""
    documents: List[Document] = []

    for file_path in file_paths:
        path = Path(file_path)
        if not path.is_file():
            continue

        documents.extend(_load_file(path))

    return documents
